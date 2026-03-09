"""
generate_pptx.py
Generates RTAI_LinkedIn_Presentation.pptx — a 5-slide dark-themed
professional deck summarising the RTAI framework and Lab_Machine_138 findings.

Usage:
    python generate_pptx.py
"""
from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
BG       = RGBColor(0x0E, 0x11, 0x17)   # #0E1117  dark background
CARD     = RGBColor(0x1A, 0x1D, 0x2E)   # #1A1D2E  card / panel
ACCENT   = RGBColor(0xFF, 0x4B, 0x4B)   # #FF4B4B  red accent
WHITE    = RGBColor(0xE0, 0xE0, 0xE0)   # #E0E0E0  body text
MUTED    = RGBColor(0x7B, 0x80, 0x99)   # #7B8099  subtext
C_CRIT   = RGBColor(0xFF, 0x3B, 0x3B)   # #FF3B3B  Critical
C_HIGH   = RGBColor(0xFF, 0x8C, 0x00)   # #FF8C00  High
C_MED    = RGBColor(0xFF, 0xD7, 0x00)   # #FFD700  Medium
C_LOW    = RGBColor(0x4C, 0xAF, 0x50)   # #4CAF50  Low
C_GREEN  = RGBColor(0x98, 0xC3, 0x79)   # #98C379  code green

SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)

BLANK_LAYOUT_IDX = 6   # index of the fully-blank slide layout


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation) -> Any:
    layout = prs.slide_layouts[BLANK_LAYOUT_IDX]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor = BG) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide: Any, x: float, y: float, w: float, h: float,
        fill: RGBColor | None = None,
        border: RGBColor | None = None,
        border_pt: float = 1.0) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid() if fill else shape.fill.background()
    if fill:
        shape.fill.fore_color.rgb = fill
    line = shape.line
    if border:
        line.color.rgb = border
        line.width = Pt(border_pt)
    else:
        line.fill.background()
    return shape


def txt(slide: Any, text: str, x: float, y: float, w: float, h: float,
        size: float = 18, bold: bool = False, color: RGBColor = WHITE,
        align: Any = PP_ALIGN.LEFT, italic: bool = False,
        font: str = "Courier New") -> Any:
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def accent_bar(slide: Any, y: float = 0.18, h: float = 0.06) -> None:
    """Thin red horizontal rule near the top."""
    box(slide, 0.5, y, 12.33, h, fill=ACCENT, border=None)


def slide_number(slide: Any, n: int) -> None:
    txt(slide, f"{n} / 5", 12.5, 7.1, 0.7, 0.3,
        size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide1_title(prs):
    """Slide 1: Title."""
    s = blank_slide(prs)
    fill_bg(s)

    # Background grid lines (decorative)
    for i in range(1, 14):
        b = box(s, 0, i * 0.55, 13.33, 0.55,
                fill=RGBColor(0x12, 0x15, 0x1E), border=None)

    # Centre card
    box(s, 2.5, 1.8, 8.33, 4.0, fill=CARD,
        border=RGBColor(0x2E, 0x32, 0x50), border_pt=1.5)
    accent_bar(s, y=1.85, h=0.08)

    txt(s, "🛡️  RTAI",
        2.7, 2.0, 8.0, 1.1,
        size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
        font="Courier New")

    txt(s, "Autonomous Red Team AI",
        2.7, 3.2, 8.0, 0.7,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Powered by LangGraph  ·  OpenAI  ·  Tavily  ·  Streamlit",
        2.7, 4.0, 8.0, 0.5,
        size=13, color=MUTED, align=PP_ALIGN.CENTER)

    txt(s, "From zero to full CVE-grounded pentest report — autonomously.",
        2.7, 4.6, 8.0, 0.5,
        size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    slide_number(s, 1)


def slide2_pipeline(prs):
    """Slide 2: 5-Stage Agentic Pipeline."""
    s = blank_slide(prs)
    fill_bg(s)
    accent_bar(s)

    txt(s, "The 5-Stage Agentic Pipeline",
        0.5, 0.3, 12.0, 0.7,
        size=26, bold=True, color=ACCENT)

    stages = [
        ("01", "ReconAgent",        "Nmap scan\nservice / OS detection",     ACCENT),
        ("02", "OsintAgent",        "Tavily CVE search\ntop-3 risk synthesis", RGBColor(0xFF,0x8C,0x00)),
        ("03", "ExploitAgent",      "Attack vector ranking\nCVSS-grounded",   RGBColor(0xFF,0xD7,0x00)),
        ("04", "RemediationAgent",  "Per-vector steps\ncode + verification",  C_LOW),
        ("05", "ReportAgent",       "Markdown report\n+ CISO Dashboard",      RGBColor(0x61,0xAF,0xEF)),
    ]

    card_w = 2.2
    gap    = 0.26
    start_x = 0.5

    for i, (num, name, desc, color) in enumerate(stages):
        x = start_x + i * (card_w + gap)

        # Card
        box(s, x, 1.2, card_w, 4.8, fill=CARD,
            border=color, border_pt=2.0)

        # Number badge
        box(s, x + 0.1, 1.3, 0.55, 0.45, fill=color, border=None)
        txt(s, num, x + 0.1, 1.3, 0.55, 0.45,
            size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)

        # Agent name
        txt(s, name, x + 0.08, 1.85, card_w - 0.16, 0.7,
            size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

        # Description
        txt(s, desc, x + 0.08, 2.65, card_w - 0.16, 1.2,
            size=11, color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow (except last)
        if i < len(stages) - 1:
            ax = x + card_w + 0.02
            txt(s, "▶", ax, 3.1, gap + 0.05, 0.4,
                size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # State label at bottom
    txt(s, "All agents share a single typed RTAIState object — findings accumulate across nodes.",
        0.5, 6.6, 12.3, 0.5,
        size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    slide_number(s, 2)


def slide3_findings(prs):
    """Slide 3: Critical Findings Summary."""
    s = blank_slide(prs)
    fill_bg(s)
    accent_bar(s)

    txt(s, "Critical Findings — Lab Machine 10.x.x.x",
        0.5, 0.3, 12.0, 0.7,
        size=26, bold=True, color=ACCENT)

    txt(s, "Target: intentionally vulnerable lab machine  ·  5 open ports  ·  3 High-risk CVEs discovered autonomously",
        0.5, 1.0, 12.3, 0.4,
        size=12, color=MUTED, italic=True)

    findings = [
        ("HIGH",     C_HIGH,  "CVE-2017-7494  —  SambaCry",
         "Samba 3.x–4.x  ·  Ports 139, 445",
         "Remote Code Execution via malicious shared library upload. "
         "CVSS 7.5. Public Metasploit module available."),

        ("HIGH",     C_HIGH,  "micro_httpd  —  Pre-Auth RCE",
         "micro_httpd  ·  Ports 80, 443",
         "Stack-based buffer overflow allowing unauthenticated remote "
         "code execution. No CVE assigned; PoC published."),

        ("HIGH",     C_HIGH,  "CVE-2021-41773  —  Apache Path Traversal",
         "Apache HTTP Server 2.4.49/2.4.50  ·  Port 80",
         "Unauthenticated path traversal enabling remote file disclosure. "
         "CVSS 7.5. Actively exploited in the wild."),
    ]

    for i, (risk, color, title, service, desc) in enumerate(findings):
        y = 1.55 + i * 1.65

        # Card
        box(s, 0.5, y, 12.33, 1.45, fill=CARD,
            border=color, border_pt=1.5)

        # Risk badge
        box(s, 0.55, y + 0.1, 0.95, 0.38, fill=color, border=None)
        txt(s, risk, 0.55, y + 0.1, 0.95, 0.38,
            size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)

        # Title
        txt(s, title, 1.65, y + 0.05, 10.9, 0.45,
            size=15, bold=True, color=WHITE)

        # Service
        txt(s, service, 1.65, y + 0.5, 10.9, 0.3,
            size=11, color=color)

        # Description
        txt(s, desc, 1.65, y + 0.82, 10.9, 0.5,
            size=10.5, color=MUTED, italic=True)

    slide_number(s, 3)


def slide4_remediation(prs):
    """Slide 4: Automated Remediation Example."""
    s = blank_slide(prs)
    fill_bg(s)
    accent_bar(s)

    txt(s, "Automated Remediation — Zero Human Input",
        0.5, 0.3, 12.0, 0.7,
        size=26, bold=True, color=ACCENT)

    txt(s, "RemediationAgent generates structured, copy-paste ready fixes for every finding.",
        0.5, 1.0, 12.3, 0.4,
        size=12, color=MUTED, italic=True)

    # Left column — finding card
    box(s, 0.5, 1.5, 5.5, 5.4, fill=CARD,
        border=C_HIGH, border_pt=1.5)
    txt(s, "FINDING",      0.6, 1.55, 2.0, 0.35, size=10, color=MUTED, bold=True)
    txt(s, "CVE-2017-7494",0.6, 1.9,  5.2, 0.5,  size=17, bold=True, color=C_HIGH)
    txt(s, "SambaCry — Remote Code Execution\nSamba smbd 3.X–4.X  ·  Ports 139, 445",
        0.6, 2.45, 5.2, 0.8, size=12, color=WHITE)
    txt(s, "CVSS: 7.5  ·  Risk: HIGH",
        0.6, 3.3, 5.2, 0.4, size=12, color=C_HIGH, bold=True)
    txt(s, "Public Metasploit module exists.\nExploit uploads a shared library via\nwritable SMB share.",
        0.6, 3.75, 5.2, 1.1, size=11, color=MUTED, italic=True)

    # Arrow
    txt(s, "▶", 6.1, 3.8, 0.6, 0.5, size=22, color=ACCENT, align=PP_ALIGN.CENTER)

    # Right column — remediation card
    box(s, 6.8, 1.5, 6.0, 5.4, fill=CARD,
        border=C_LOW, border_pt=1.5)
    txt(s, "REMEDIATION",  6.9, 1.55, 3.0, 0.35, size=10, color=MUTED, bold=True)
    txt(s, "Patch Samba CVE-2017-7494", 6.9, 1.9, 5.7, 0.5, size=15, bold=True, color=C_LOW)

    steps = [
        "1.  systemctl stop smb",
        "2.  cp /etc/samba/smb.conf smb.conf.bak",
        "3.  apt-get update",
        "4.  apt-get install --only-upgrade samba",
        "5.  systemctl start smb",
    ]
    box(s, 6.9, 2.5, 5.7, 2.3, fill=BG,
        border=RGBColor(0x2E,0x32,0x50), border_pt=1.0)
    txt(s, "\n".join(steps), 7.0, 2.55, 5.5, 2.2,
        size=11, color=C_GREEN, font="Courier New")

    txt(s, "VERIFICATION", 6.9, 4.9, 2.5, 0.3, size=10, color=MUTED, bold=True)
    box(s, 6.9, 5.25, 5.7, 0.45, fill=BG,
        border=RGBColor(0x2E,0x32,0x50), border_pt=1.0)
    txt(s, "smbd -V | grep '4\\.6\\.5\\|4\\.7\\.0'",
        7.0, 5.28, 5.5, 0.38,
        size=11, color=C_GREEN, font="Courier New")

    slide_number(s, 4)


def slide5_roadmap(prs):
    """Slide 5: Future Roadmap & Dashboard."""
    s = blank_slide(prs)
    fill_bg(s)
    accent_bar(s)

    txt(s, "Roadmap & CISO Dashboard",
        0.5, 0.3, 12.0, 0.7,
        size=26, bold=True, color=ACCENT)

    # Left: dashboard feature list
    box(s, 0.5, 1.2, 5.8, 5.6, fill=CARD,
        border=RGBColor(0x2E,0x32,0x50), border_pt=1.5)
    txt(s, "🛡️  CISO Dashboard (Live)",
        0.7, 1.3, 5.4, 0.5, size=16, bold=True, color=ACCENT)

    dashboard_items = [
        ("✅", "Dark-themed Streamlit UI"),
        ("✅", "Plotly grouped bar chart — risk by engagement"),
        ("✅", "Donut chart — per-engagement breakdown"),
        ("✅", "Metric cards: Total / Critical / High / Medium / Low"),
        ("✅", "Open port chips per engagement"),
        ("✅", "Full report viewer with rendered Markdown"),
    ]
    for i, (icon, label) in enumerate(dashboard_items):
        txt(s, f"{icon}  {label}",
            0.7, 1.9 + i * 0.55, 5.4, 0.48,
            size=12, color=WHITE)

    # Right: roadmap
    box(s, 6.83, 1.2, 6.0, 5.6, fill=CARD,
        border=RGBColor(0x2E,0x32,0x50), border_pt=1.5)
    txt(s, "🗺️  Next Milestones",
        7.03, 1.3, 5.6, 0.5, size=16, bold=True, color=RGBColor(0x61,0xAF,0xEF))

    roadmap = [
        (C_CRIT, "v2.0", "CISO Report Export to PDF"),
        (C_HIGH, "v2.1", "Nuclei / Metasploit Tool Wrappers"),
        (C_HIGH, "v2.2", "Shadow Agent — IDS-evasive scanning"),
        (C_MED,  "v2.3", "Multi-target CIDR sweep mode"),
        (C_MED,  "v2.4", "Slack / email alert integration"),
        (C_LOW,  "v3.0", "Real-time collaborative dashboard"),
    ]
    for i, (color, version, label) in enumerate(roadmap):
        y = 1.9 + i * 0.75
        box(s, 7.03, y, 0.7, 0.38, fill=color, border=None)
        txt(s, version, 7.03, y, 0.7, 0.38,
            size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)
        txt(s, label, 7.85, y, 4.8, 0.38, size=12, color=WHITE)

    slide_number(s, 5)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()
    slide1_title(prs)
    slide2_pipeline(prs)
    slide3_findings(prs)
    slide4_remediation(prs)
    slide5_roadmap(prs)

    out = "RTAI_LinkedIn_Presentation.pptx"
    prs.save(out)
    print(f"[+] Saved: {out}")


if __name__ == "__main__":
    main()
